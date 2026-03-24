"""Build a QC report PowerPoint directly from parsed QC markdown output."""
from __future__ import annotations

import re
from pathlib import Path

from pydantic import BaseModel
from pptx import Presentation
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN
from pptx.util import Inches, Pt, Emu

# ── Colors (from build_ppt.yaml design spec) ────────────────────────────
NAVY = RGBColor(0x1B, 0x2A, 0x4A)
TEAL = RGBColor(0x00, 0xA4, 0x99)
WHITE = RGBColor(0xFF, 0xFF, 0xFF)
LIGHT_GRAY = RGBColor(0xF5, 0xF5, 0xF5)
DARK_GRAY = RGBColor(0x61, 0x61, 0x61)
BAND_ALT = RGBColor(0xFA, 0xFA, 0xFA)

RATING_STYLE = {
    "MEETS":           {"bg": RGBColor(0xE8, 0xF5, 0xE9), "fg": RGBColor(0x2E, 0x7D, 0x32)},
    "PARTIALLY MEETS": {"bg": RGBColor(0xFF, 0xF3, 0xE0), "fg": RGBColor(0xE6, 0x51, 0x00)},
    "DOES NOT MEET":   {"bg": RGBColor(0xFF, 0xEB, 0xEE), "fg": RGBColor(0xC6, 0x28, 0x28)},
}

PRIORITY_COLOR = {
    "HIGH":   RGBColor(0xC6, 0x28, 0x28),
    "MEDIUM": RGBColor(0xE6, 0x51, 0x00),
    "LOW":    RGBColor(0x61, 0x61, 0x61),
}

SLIDE_W = Inches(13.333)
SLIDE_H = Inches(7.5)


# ── Parsed data models ──────────────────────────────────────────────────
class QCRow(BaseModel):
    criteria: str
    rating: str
    rationale: str
    revision: str = "N/A"


class OverallRow(BaseModel):
    label: str
    rating: str
    rationale: str


class QCReport(BaseModel):
    qc_rows: list[QCRow] = []
    overall_rows: list[OverallRow] = []
    revised_description: str = ""


# ── Markdown parser ─────────────────────────────────────────────────────
def _split_sections(text: str) -> dict[str, str]:
    """Split QC output by ### headers into {header_lower: body}."""
    parts: dict[str, str] = {}
    current_key = ""
    lines: list[str] = []
    for line in text.split("\n"):
        m = re.match(r"^###?\s*\d*\.?\s*(.+)", line)
        if m:
            if current_key:
                parts[current_key] = "\n".join(lines).strip()
            current_key = m.group(1).strip().lower()
            lines = []
        else:
            lines.append(line)
    if current_key:
        parts[current_key] = "\n".join(lines).strip()
    return parts


def _parse_md_table(text: str, ncols: int) -> list[list[str]]:
    """Parse a markdown table, returning rows (excluding header + separator)."""
    rows: list[list[str]] = []
    for line in text.split("\n"):
        line = line.strip()
        if not line.startswith("|") or re.match(r"^\|[\s\-:|]+\|$", line):
            continue
        # strip leading/trailing pipes, then split on inner pipes
        cells = [c.strip() for c in line.strip("|").split("|")]
        if len(cells) == ncols:
            rows.append(cells)
    # first row is header; skip it
    return rows[1:] if len(rows) > 1 else []


def parse_qc_output(text: str) -> QCReport:
    """Parse LLM QC markdown output into structured QCReport."""
    sections = _split_sections(text)

    # QC table (4-col: Criteria | PASS/FAIL | Rationale | Revision)
    qc_rows: list[QCRow] = []
    for key in sections:
        if "quality control" in key or "qc" in key.split():
            for row in _parse_md_table(sections[key], 4):
                qc_rows.append(QCRow(criteria=row[0], rating=row[1], rationale=row[2], revision=row[3]))
            break

    # Overall assessment (3-col: Assessment/QC Type | PASS/FAIL | Rationale)
    overall_rows: list[OverallRow] = []
    for key in sections:
        if "overall" in key:
            for row in _parse_md_table(sections[key], 3):
                overall_rows.append(OverallRow(label=row[0], rating=row[1], rationale=row[2]))
            break

    # Revised description (plain text after header)
    revised = ""
    for key in sections:
        if "revised" in key:
            revised = sections[key].strip()
            break

    return QCReport(qc_rows=qc_rows, overall_rows=overall_rows, revised_description=revised)


# ── PPTX helpers ────────────────────────────────────────────────────────
def _set_bg(slide, color: RGBColor):
    fill = slide.background.fill
    fill.solid()
    fill.fore_color.rgb = color


def _add_textbox(slide, left, top, width, height, text, *,
                 font_size=12, bold=False, color=NAVY, alignment=PP_ALIGN.LEFT, word_wrap=True):
    txbox = slide.shapes.add_textbox(left, top, width, height)
    tf = txbox.text_frame
    tf.word_wrap = word_wrap
    p = tf.paragraphs[0]
    p.text = text
    p.font.size = Pt(font_size)
    p.font.bold = bold
    p.font.color.rgb = color
    p.alignment = alignment
    return txbox


def _style_cell(cell, text, *, font_size=9, bold=False, color=NAVY, bg=None):
    cell.text = ""
    p = cell.text_frame.paragraphs[0]
    p.text = text
    p.font.size = Pt(font_size)
    p.font.bold = bold
    p.font.color.rgb = color
    cell.text_frame.word_wrap = True
    if bg:
        cell.fill.solid()
        cell.fill.fore_color.rgb = bg


def _rating_cell(cell, rating: str):
    """Color-code a rating cell."""
    key = rating.strip().upper()
    style = RATING_STYLE.get(key, RATING_STYLE["DOES NOT MEET"])
    _style_cell(cell, rating.strip(), bold=True, color=style["fg"], bg=style["bg"], font_size=9)


def _add_colored_rect(slide, left, top, width, height, color: RGBColor):
    """Add a thin colored rectangle (used as a border accent)."""
    shape = slide.shapes.add_shape(1, left, top, width, height)  # MSO_SHAPE.RECTANGLE = 1
    shape.fill.solid()
    shape.fill.fore_color.rgb = color
    shape.line.fill.background()


# ── Slide builders ──────────────────────────────────────────────────────
def _slide_title(prs: Presentation, report: QCReport, original_control: str):
    slide = prs.slides.add_slide(prs.slide_layouts[6])  # blank
    _set_bg(slide, NAVY)
    _add_textbox(slide, Inches(1), Inches(1.8), Inches(11), Inches(1.2),
                 "RCSA Control Description — Quality Review",
                 font_size=36, bold=True, color=WHITE, alignment=PP_ALIGN.CENTER)

    subtitle = original_control[:80] + ("…" if len(original_control) > 80 else "")
    _add_textbox(slide, Inches(1), Inches(3.2), Inches(11), Inches(0.6),
                 subtitle, font_size=18, color=LIGHT_GRAY, alignment=PP_ALIGN.CENTER)

    # badge row
    badges = ["QC Assessment", "Overall Rating", "Revised Description"]
    for i, label in enumerate(badges):
        x = Inches(3.5 + i * 2.3)
        shape = slide.shapes.add_shape(1, x, Inches(4.2), Inches(2), Inches(0.4))
        shape.fill.solid()
        shape.fill.fore_color.rgb = TEAL
        shape.line.fill.background()
        tf = shape.text_frame
        tf.paragraphs[0].text = label
        tf.paragraphs[0].font.size = Pt(10)
        tf.paragraphs[0].font.color.rgb = WHITE
        tf.paragraphs[0].font.bold = True
        tf.paragraphs[0].alignment = PP_ALIGN.CENTER

    _add_textbox(slide, Inches(1), Inches(6.5), Inches(11), Inches(0.4),
                 "Internal Audit · Quality Control · Confidential",
                 font_size=9, color=LIGHT_GRAY, alignment=PP_ALIGN.CENTER)


def _slide_context(prs: Presentation, report: QCReport, original_control: str):
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    _add_textbox(slide, Inches(0.5), Inches(0.3), Inches(12), Inches(0.6),
                 "Evaluation Context", font_size=24, bold=True, color=NAVY)

    # Left: QC Framework
    _add_colored_rect(slide, Inches(0.5), Inches(1.2), Inches(5.5), Inches(0.45), NAVY)
    _add_textbox(slide, Inches(0.6), Inches(1.2), Inches(5.3), Inches(0.45),
                 "QC Framework", font_size=14, bold=True, color=WHITE)
    criteria = [r.criteria for r in report.qc_rows] or ["Who", "What", "When", "How", "Why", "Evidence"]
    _add_textbox(slide, Inches(0.7), Inches(1.9), Inches(5), Inches(4),
                 "QC Criteria:  " + " · ".join(criteria),
                 font_size=11, color=NAVY)

    # Right: Original control
    _add_colored_rect(slide, Inches(6.5), Inches(1.2), Inches(0.06), Inches(5.5), TEAL)
    _add_textbox(slide, Inches(6.3), Inches(1.2), Inches(6.5), Inches(0.3),
                 "ORIGINAL CONTROL DESCRIPTION", font_size=9, bold=True, color=DARK_GRAY)
    _add_textbox(slide, Inches(6.8), Inches(1.7), Inches(6), Inches(5),
                 original_control, font_size=11, color=NAVY)


def _slide_qc_table(prs: Presentation, report: QCReport):
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    met = sum(1 for r in report.qc_rows if r.rating.strip().upper() == "MEETS")
    total = len(report.qc_rows)

    _add_textbox(slide, Inches(0.5), Inches(0.3), Inches(10), Inches(0.5),
                 "Control Description Quality Control Report",
                 font_size=24, bold=True, color=NAVY)
    _add_textbox(slide, Inches(0.5), Inches(0.85), Inches(6), Inches(0.35),
                 f"{met} of {total} criteria met",
                 font_size=12, color=DARK_GRAY)

    if not report.qc_rows:
        return

    rows_n = total + 1  # +header
    cols_n = 4
    tbl = slide.shapes.add_table(rows_n, cols_n, Inches(0.5), Inches(1.4), Inches(12.3), Inches(5.5)).table

    # column widths
    tbl.columns[0].width = Inches(1.7)
    tbl.columns[1].width = Inches(1.7)
    tbl.columns[2].width = Inches(5.2)
    tbl.columns[3].width = Inches(3.7)

    for ci, header in enumerate(["QC Criteria", "PASS/FAIL", "Rationale", "Revision"]):
        _style_cell(tbl.cell(0, ci), header, font_size=10, bold=True, color=WHITE, bg=NAVY)

    for ri, row in enumerate(report.qc_rows, start=1):
        bg = WHITE if ri % 2 else BAND_ALT
        _style_cell(tbl.cell(ri, 0), row.criteria, bg=bg)
        _rating_cell(tbl.cell(ri, 1), row.rating)
        _style_cell(tbl.cell(ri, 2), row.rationale, bg=bg)
        _style_cell(tbl.cell(ri, 3), row.revision, bg=bg)


def _slide_assessment(prs: Presentation, report: QCReport):
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    _add_textbox(slide, Inches(0.5), Inches(0.3), Inches(10), Inches(0.5),
                 "Overall Assessment", font_size=24, bold=True, color=NAVY)

    if not report.overall_rows:
        _add_textbox(slide, Inches(0.5), Inches(2), Inches(12), Inches(1),
                     "No overall assessment data found.", font_size=14, color=DARK_GRAY)
        return

    card_w = Inches(3.8)
    card_h = Inches(4.5)
    gap = Inches(0.5)
    start_x = Inches(0.5)

    for i, row in enumerate(report.overall_rows):
        x = start_x + i * (card_w + gap)
        y = Inches(1.2)

        style = RATING_STYLE.get(row.rating.strip().upper(), RATING_STYLE["DOES NOT MEET"])
        # top border accent
        _add_colored_rect(slide, x, y, card_w, Inches(0.06), style["fg"])
        # card bg
        card = slide.shapes.add_shape(1, x, y + Inches(0.06), card_w, card_h)
        card.fill.solid()
        card.fill.fore_color.rgb = WHITE
        card.line.color.rgb = LIGHT_GRAY

        _add_textbox(slide, x + Inches(0.3), y + Inches(0.2), card_w - Inches(0.6), Inches(0.4),
                     row.label, font_size=12, bold=True, color=NAVY)
        _add_textbox(slide, x + Inches(0.3), y + Inches(0.7), card_w - Inches(0.6), Inches(0.6),
                     row.rating.strip(), font_size=20, bold=True, color=style["fg"],
                     alignment=PP_ALIGN.CENTER)
        _add_textbox(slide, x + Inches(0.3), y + Inches(1.5), card_w - Inches(0.6), card_h - Inches(1.8),
                     row.rationale, font_size=9, color=DARK_GRAY)


def _slide_revision(prs: Presentation, report: QCReport, original_control: str):
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    _add_textbox(slide, Inches(0.5), Inches(0.3), Inches(12), Inches(0.5),
                 "Recommended Control Description Revision",
                 font_size=24, bold=True, color=NAVY)

    panel_w = Inches(5.8)
    panel_h = Inches(5.5)

    # Before panel
    _add_colored_rect(slide, Inches(0.5), Inches(1.2), Inches(0.06), panel_h, RGBColor(0xC6, 0x28, 0x28))
    _add_textbox(slide, Inches(0.7), Inches(1.2), panel_w, Inches(0.3),
                 "ORIGINAL CONTROL DESCRIPTION", font_size=9, bold=True,
                 color=RGBColor(0xC6, 0x28, 0x28))
    _add_textbox(slide, Inches(0.7), Inches(1.7), panel_w, panel_h - Inches(0.5),
                 original_control, font_size=10, color=NAVY)

    # After panel
    _add_colored_rect(slide, Inches(7), Inches(1.2), Inches(0.06), panel_h, RGBColor(0x2E, 0x7D, 0x32))
    _add_textbox(slide, Inches(7.2), Inches(1.2), panel_w, Inches(0.3),
                 "FULLY REVISED CONTROL DESCRIPTION", font_size=9, bold=True,
                 color=RGBColor(0x2E, 0x7D, 0x32))
    _add_textbox(slide, Inches(7.2), Inches(1.7), panel_w, panel_h - Inches(0.5),
                 report.revised_description or "[No revised description generated]",
                 font_size=10, color=NAVY)

    # Placeholder callout
    placeholders = report.revised_description.count("[PLACEHOLDER]")
    if placeholders:
        _add_textbox(slide, Inches(0.5), Inches(6.8), Inches(12), Inches(0.4),
                     f"⚠ Action required: {placeholders} placeholder(s) must be resolved by the control owner.",
                     font_size=9, bold=True, color=RGBColor(0xE6, 0x51, 0x00))


def _slide_next_steps(prs: Presentation, report: QCReport):
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    _add_textbox(slide, Inches(0.5), Inches(0.3), Inches(12), Inches(0.5),
                 "Next Steps & Action Items", font_size=24, bold=True, color=NAVY)

    actions: list[tuple[str, str, str]] = []  # (action, owner, priority)
    for row in report.qc_rows:
        r = row.rating.strip().upper()
        if r == "DOES NOT MEET":
            actions.append((f"Revise control to address {row.criteria}", "Control Owner", "HIGH"))
        elif r == "PARTIALLY MEETS":
            actions.append((f"Strengthen {row.criteria} language", "Control Owner", "MEDIUM"))
    if "[PLACEHOLDER]" in report.revised_description:
        actions.append(("Supply missing placeholder details", "Control Owner", "HIGH"))
    actions.append(("Review revised description against RCSA system of record", "Risk Manager", "MEDIUM"))
    actions.append(("Re-run automated QC validation after revision", "Internal Audit", "LOW"))

    rows_n = len(actions) + 1
    tbl = slide.shapes.add_table(rows_n, 3, Inches(0.5), Inches(1.2), Inches(12.3), Inches(0.5 * rows_n)).table
    tbl.columns[0].width = Inches(7.5)
    tbl.columns[1].width = Inches(2.4)
    tbl.columns[2].width = Inches(2.4)

    for ci, h in enumerate(["Action", "Owner", "Priority"]):
        _style_cell(tbl.cell(0, ci), h, font_size=10, bold=True, color=WHITE, bg=NAVY)

    for ri, (action, owner, priority) in enumerate(actions, start=1):
        bg = WHITE if ri % 2 else BAND_ALT
        _style_cell(tbl.cell(ri, 0), action, bg=bg)
        _style_cell(tbl.cell(ri, 1), owner, bg=bg)
        _style_cell(tbl.cell(ri, 2), priority, bold=True,
                    color=PRIORITY_COLOR.get(priority, DARK_GRAY), bg=bg)

    _add_textbox(slide, Inches(0.5), Inches(6.5), Inches(12), Inches(0.4),
                 "All HIGH-priority actions must be completed before the next RCSA cycle.",
                 font_size=9, bold=True, color=NAVY)


# ── Public API ──────────────────────────────────────────────────────────
def build_qc_pptx(original_control: str, qc_output: str, output_path: Path) -> Path:
    """Parse QC markdown and generate a 6-slide PPTX deck."""
    report = parse_qc_output(qc_output)

    prs = Presentation()
    prs.slide_width = SLIDE_W
    prs.slide_height = SLIDE_H

    _slide_title(prs, report, original_control)
    _slide_context(prs, report, original_control)
    _slide_qc_table(prs, report)
    _slide_assessment(prs, report)
    _slide_revision(prs, report, original_control)
    _slide_next_steps(prs, report)

    output_path = Path(output_path)
    output_path.parent.mkdir(parents=True, exist_ok=True)
    prs.save(str(output_path))
    return output_path
